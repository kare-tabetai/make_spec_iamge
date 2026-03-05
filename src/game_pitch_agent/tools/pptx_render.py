"""python-pptx によるペラ1企画書 PPTX 生成モジュール"""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

logger = logging.getLogger(__name__)

# ── カラーパレット ──
COLOR_TITLE_BG = RGBColor(0x1B, 0x2A, 0x4A)   # 濃紺
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ACCENT = RGBColor(0x2E, 0x50, 0x90)       # アクセント青
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)          # ダークグレイ
COLOR_LIGHT_GRAY = RGBColor(0xDD, 0xDD, 0xDD)   # プレースホルダー背景
COLOR_BADGE_BG = RGBColor(0x3A, 0x60, 0xA0)     # バッジ背景
COLOR_SECTION_BG = RGBColor(0xF5, 0xF7, 0xFA)   # セクション背景

FONT_NAME = "Meiryo"

# 16:9 スライドサイズ (Emu)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# レイアウト定数
MARGIN = Inches(0.3)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN * 2
COL_WIDTH = (CONTENT_WIDTH - Inches(0.2)) / 2  # 2カラム用


def render_pitch_pptx(
    pitch: dict,
    output_path: str,
    image_path: str | None = None,
) -> str:
    """企画書データから PPTX ファイルを生成する。

    Args:
        pitch: ExpandedPitch フィールドの辞書
        output_path: 保存先パス (.pptx)
        image_path: 埋め込み画像パス (あれば)

    Returns:
        保存先パス文字列
    """
    prs = _create_presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    y_cursor = _add_title_bar(slide, pitch)
    y_cursor = _add_two_column_row(
        slide, y_cursor,
        left_title="コンセプト",
        left_body=pitch.get("concept", ""),
        right_title="コアメカニクス",
        right_body=pitch.get("core_mechanic", ""),
    )
    y_cursor = _add_two_column_row(
        slide, y_cursor,
        left_title="ゲーム概要",
        left_body=pitch.get("overview", ""),
        right_title="ゲームサイクル",
        right_body=None,
        right_custom=_game_cycle_text(pitch),
    )
    _add_two_column_row(
        slide, y_cursor,
        left_title="アートスタイル",
        left_body=pitch.get("art_style", ""),
        left_image_path=image_path,
        right_title="差別化ポイント（USP）",
        right_body=pitch.get("usp", ""),
        right_extra_title="実現可能性",
        right_extra_body=pitch.get("feasibility_note", ""),
    )

    prs.save(output_path)
    logger.info(f"PPTX保存完了: {output_path}")
    return output_path


# ── 内部ヘルパー ──


def _create_presentation() -> Presentation:
    """16:9 空のプレゼンテーションを作成する。"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _add_textbox(
    slide,
    left, top, width, height,
    text: str,
    font_size: int = 11,
    bold: bool = False,
    color: RGBColor = COLOR_BODY,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = FONT_NAME,
) -> None:
    """汎用テキストボックスを追加する。"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment


def _add_rect(slide, left, top, width, height, fill_color: RGBColor):
    """塗りつぶし矩形を追加する。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_title_bar(slide, pitch: dict) -> int:
    """タイトルバーを描画し、次のY座標を返す。"""
    title = pitch.get("title", "無題")
    catchcopy = pitch.get("catchcopy", "")
    genre = pitch.get("genre", "")
    platform = pitch.get("platform", "")

    bar_height = Inches(1.1)

    # 背景矩形
    _add_rect(slide, 0, 0, SLIDE_WIDTH, bar_height, COLOR_TITLE_BG)

    # タイトル
    _add_textbox(
        slide,
        left=MARGIN, top=Inches(0.15),
        width=CONTENT_WIDTH * 0.7, height=Inches(0.5),
        text=title,
        font_size=28, bold=True, color=COLOR_WHITE,
    )

    # キャッチコピー
    if catchcopy:
        _add_textbox(
            slide,
            left=MARGIN, top=Inches(0.65),
            width=CONTENT_WIDTH * 0.7, height=Inches(0.35),
            text=catchcopy,
            font_size=14, color=COLOR_WHITE,
        )

    # ジャンル / プラットフォーム バッジ
    badge_text = " / ".join(filter(None, [genre, platform]))
    if badge_text:
        badge_left = SLIDE_WIDTH - MARGIN - Inches(3.5)
        badge = _add_rect(
            slide, badge_left, Inches(0.25),
            Inches(3.2), Inches(0.45),
            COLOR_BADGE_BG,
        )
        badge.text_frame.word_wrap = True
        p = badge.text_frame.paragraphs[0]
        p.text = badge_text
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER

    return bar_height


def _add_section_box(
    slide,
    left, top, width, height,
    title: str,
    body: str,
    image_path: str | None = None,
    extra_title: str | None = None,
    extra_body: str | None = None,
) -> None:
    """セクション（ヘッダー + 本文）を描画する。"""
    padding = Inches(0.12)
    header_h = Inches(0.35)

    # セクション背景
    _add_rect(slide, left, top, width, height, COLOR_SECTION_BG)

    # セクションヘッダー
    _add_rect(slide, left, top, width, header_h, COLOR_ACCENT)
    _add_textbox(
        slide,
        left=left + padding, top=top + Inches(0.04),
        width=width - padding * 2, height=header_h,
        text=title,
        font_size=12, bold=True, color=COLOR_WHITE,
    )

    body_top = top + header_h + Inches(0.08)
    body_width = width - padding * 2

    # 画像 or プレースホルダー
    if image_path and Path(image_path).exists():
        img_height = Inches(1.5)
        slide.shapes.add_picture(
            image_path,
            left + padding, body_top,
            body_width, img_height,
        )
        body_top += img_height + Inches(0.08)
    elif image_path is not None:
        # プレースホルダー矩形
        ph_height = Inches(1.2)
        ph = _add_rect(slide, left + padding, body_top, body_width, ph_height, COLOR_LIGHT_GRAY)
        ph.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        ph.line.dash_style = 2  # dash
        ph_tf = ph.text_frame
        ph_tf.paragraphs[0].text = "Image Placeholder"
        ph_tf.paragraphs[0].font.size = Pt(10)
        ph_tf.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        ph_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        ph_tf.paragraphs[0].font.name = FONT_NAME
        body_top += ph_height + Inches(0.08)

    # 本文
    if body:
        remaining = top + height - body_top - Inches(0.05)
        if extra_title and extra_body:
            remaining = remaining // 2
        _add_textbox(
            slide,
            left=left + padding, top=body_top,
            width=body_width, height=max(remaining, Inches(0.3)),
            text=body,
            font_size=10,
        )
        body_top += remaining

    # 追加セクション（USP + 実現可能性 の2段構成用）
    if extra_title and extra_body:
        extra_header_h = Inches(0.3)
        _add_rect(slide, left, body_top + Inches(0.05), width, extra_header_h, COLOR_ACCENT)
        _add_textbox(
            slide,
            left=left + padding, top=body_top + Inches(0.08),
            width=body_width, height=extra_header_h,
            text=extra_title,
            font_size=11, bold=True, color=COLOR_WHITE,
        )
        _add_textbox(
            slide,
            left=left + padding, top=body_top + extra_header_h + Inches(0.1),
            width=body_width, height=top + height - body_top - extra_header_h - Inches(0.15),
            text=extra_body,
            font_size=10,
        )


def _add_two_column_row(
    slide,
    y_top: int,
    left_title: str,
    left_body: str | None,
    right_title: str,
    right_body: str | None,
    right_custom: str | None = None,
    left_image_path: str | None = None,
    right_extra_title: str | None = None,
    right_extra_body: str | None = None,
) -> int:
    """2カラムの行を描画し、次のY座標を返す。"""
    row_height = Inches(1.9)
    gap = Inches(0.2)

    _add_section_box(
        slide,
        left=MARGIN, top=y_top + Inches(0.1),
        width=COL_WIDTH, height=row_height,
        title=left_title,
        body=left_body or "",
        image_path=left_image_path,
    )

    right_body_text = right_custom if right_custom else (right_body or "")
    _add_section_box(
        slide,
        left=MARGIN + COL_WIDTH + gap, top=y_top + Inches(0.1),
        width=COL_WIDTH, height=row_height,
        title=right_title,
        body=right_body_text,
        extra_title=right_extra_title,
        extra_body=right_extra_body,
    )

    return y_top + row_height + Inches(0.15)


def _game_cycle_text(pitch: dict) -> str:
    """ゲームサイクルを箇条書きテキストに変換する。"""
    game_cycle = pitch.get("game_cycle", {})
    if not isinstance(game_cycle, dict):
        return ""
    lines = []
    main_action = game_cycle.get("main_action", "")
    short_reward = game_cycle.get("short_term_reward", "")
    long_reward = game_cycle.get("long_term_reward", "")
    if main_action:
        lines.append(f"● メインアクション: {main_action}")
    if short_reward:
        lines.append(f"● 短期的な報酬: {short_reward}")
    if long_reward:
        lines.append(f"● 中長期的な報酬: {long_reward}")
    return "\n".join(lines)
